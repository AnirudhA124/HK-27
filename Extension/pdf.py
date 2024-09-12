from flask import Flask, request, jsonify
from flask_cors import CORS
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_fireworks import ChatFireworks
from langchain_cohere import CohereEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from langchain import hub
from tempfile import NamedTemporaryFile
from pptx import Presentation
from langchain_core.documents import Document
from docx import Document as DocxDocument
import docx2txt

app = Flask(__name__)
CORS(app)

llm = ChatFireworks(api_key="0YdGG4CL6KUAgR5v207G4kBb2rWvNkXoLDbyrHxc89ag3PVt", model="accounts/fireworks/models/llama-v3-70b-instruct")
prompt = hub.pull("chaiboi/pdf_text_prompt", api_key='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')
db = None

@app.route('/')
def index():
    return "Hello world"

@app.route('/upload', methods=['POST'])
def upload():
    global db
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file:
        file_extension = file.filename.split('.')[-1].lower()
        print(file_extension)

        if file_extension in ['pdf', 'ppt', 'pptx', 'doc', 'docx', 'txt']:
            with NamedTemporaryFile(delete=False, suffix=f".{file_extension}") as temp_file:
                file.save(temp_file.name)
                temp_file_path = temp_file.name
                print(temp_file_path)

                try:
                    if file_extension == 'pdf':
                        loader = PyPDFLoader(temp_file_path)
                        docs = loader.load()
                    elif file_extension in ['ppt', 'pptx']:
                        docs = load_ppt_text(temp_file_path)
                    elif file_extension in ['doc', 'docx']:
                        docs = load_doc_text(temp_file_path)
                    elif file_extension == 'txt':
                        docs = load_txt_text(temp_file_path)
                    embeddings_model = CohereEmbeddings(cohere_api_key="uml0lVi8lxTjTL10Bkb42inOlNFk3zDf7sELxPDN",
                                                        model="embed-english-light-v3.0")
                    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1100, chunk_overlap=150)
                    splits = text_splitter.split_documents(docs)

                    db = FAISS.from_documents(splits, embeddings_model)
                    return jsonify({'fileId': temp_file.name})

                except Exception as e:
                    print(f"Error processing file: {e}")
                    return jsonify({'error': 'Error processing file'}), 500

        return jsonify({'error': 'Unsupported file type'}), 400

    return jsonify({'error': 'File upload failed'}), 500

#Function to load text from a PowerPoint file
def load_ppt_text(ppt_path):
    """Extract text from a PowerPoint file and return as document objects."""
    prs = Presentation(ppt_path)
    docs = []
    
    for idx, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        full_slide_text = "\n".join(slide_text)
        docs.append(Document(metadata={'source': ppt_path, 'page': idx}, page_content=full_slide_text))
    
    return docs

#Function to load text from a Word file
def load_doc_text(doc_path):
    """Extract text from a Word document (.docx) and return as document objects."""
    docs = []
    file_extension = doc_path.split('.')[-1].lower()
    
    if file_extension == 'doc':
        raise NotImplementedError("Handling .doc files is not implemented. Please convert to .docx.")
    
    text = docx2txt.process(doc_path)
    docs.append(Document(metadata={'source': doc_path, 'page': 0}, page_content=text))
    
    return docs

#Function to load text from a Text file
def load_txt_text(txt_path):
    docs = []
    with open(txt_path, 'r', encoding='utf-8') as file:
        text = file.read()
    docs.append(Document(metadata={'source': txt_path, 'page': 0}, page_content=text))
    return docs

def format_docs(docs):
    return "\n\n".join(doc.page_content for doc in docs)

@app.route('/chat', methods=['POST'])
def chat():
    global db
    data = request.get_json()
    user_message = data.get('message', '')
    file_id = data.get('fileId', '')

    if not db:
        return jsonify({'response': "No documents uploaded or processed."}), 400

    retriever = db.as_retriever(kwargs={"score_threshold": 0.5})
    rag_chain = (
        {"context": retriever | format_docs, "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )
    try:
        bot_response = rag_chain.invoke(user_message)
    except Exception as e:
        print(f"Error: {e}")
        bot_response = "Sorry, I encountered an error while processing your request."

    return jsonify({'response': bot_response})

if __name__ == '__main__':
    app.run(debug=True)
